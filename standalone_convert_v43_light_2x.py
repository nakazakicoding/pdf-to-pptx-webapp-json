"""
PDF to PowerPoint Converter v43 LIGHTWEIGHT (2x Resolution)
- v43ベース
- メモリ最適化版: 解像度2.0x（フル品質）、ページごとにメモリ/ディスク解放
"""
import os
import sys
import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import sys

# SSIMフォントサイズ検出モジュールをインポート（v3: 二分探索+JSON/SSIM太字判定）
from ssim_font_detector_v3 import detect_font_properties_v3, normalize_font_sizes

# OCRで検出した全ワードリスト（グローバル）
OCR_ALL_WORDS = []

# 現在処理中のページ番号（1始まり）
CURRENT_PAGE_NUM = 1

# JSONファイルパス（CLIから設定される）
CURRENT_JSON_PATH = None

# OCRbboxマッチング関数をインポート (v3: 改良版行検出)
from ocr_bbox_matcher_v4 import find_ocr_bbox_for_text



def set_font_for_run(run, font_name):
    """日本語フォントを含めてフォントを設定する"""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    
    # East Asian Font (日本語)
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'))
        rPr.append(ea)
    ea.set('typeface', font_name)
    
    # Latin Font (英数字)
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = rPr.makeelement(qn('a:latin'))
        rPr.insert(0, latin) # latinは先頭に近い方がいいかも
    latin.set('typeface', font_name)
import unicodedata
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from collections import Counter, defaultdict
import io
import json

# v27: 新JSONフォーマット（all_pages_ocr_v2.json）の読み込み関数
def load_ocr_corrections_v2():
    """新しいJSONフォーマット（v2）を読み込む"""
    global CURRENT_JSON_PATH
    
    # CLIから渡されたパスを優先
    if CURRENT_JSON_PATH and os.path.exists(CURRENT_JSON_PATH):
        json_path = CURRENT_JSON_PATH
    else:
        # フォールバック: デフォルトパス
        base_dir = os.path.dirname(os.path.abspath(__file__))
        json_path = os.path.join(base_dir, "temp_processing", "image_analysis.json")
    
    print(f"Loading JSON v2 from: {json_path}")
    if os.path.exists(json_path):
        try:
            with open(json_path, "r", encoding="utf-8") as f:
                data = json.load(f)
                print(f"Loaded keys: {list(data.keys())}")
                return data
        except Exception as e:
            print(f"Error loading corrections from {json_path}: {e}")
            return {}
    else:
        print(f"Warning: JSON not found at {json_path}")
        return {}

# 後方互換性のために旧関数も残す
def load_ocr_corrections():
    return load_ocr_corrections_v2()
import tempfile
import shutil

# 出力ディレクトリ（英語パス使用）
OUTPUT_DIR = tempfile.mkdtemp(prefix="pdf_pptx_")
print(f"Output directory: {OUTPUT_DIR}")

def read_image_cv2(image_path):
    """日本語パス対応でOpenCV画像読み込み"""
    img_array = np.fromfile(image_path, dtype=np.uint8)
    img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)
    return img

def detect_font_weight_v5(img, bbox, bg_color):
    """v5: テキスト領域の太字・明朝判定 (SSIMなしの画素解析版。既存ロジックを強化)"""
    x1, y1, x2, y2 = [int(v) for v in bbox]
    h, w = img.shape[:2]
    
    x1, x2 = max(0, x1), min(w, x2)
    y1, y2 = max(0, y1), min(h, y2)
    
    if x2 <= x1 or y2 <= y1:
        return 0.0, 0.0
    
    region = img[y1:y2, x1:x2]
    gray = cv2.cvtColor(region, cv2.COLOR_BGR2GRAY)
    
    bg_gray = cv2.cvtColor(np.array([[bg_color]], dtype=np.uint8), cv2.COLOR_BGR2GRAY)[0][0]
    if bg_gray > 128:
        _, binary = cv2.threshold(gray, bg_gray - 30, 255, cv2.THRESH_BINARY_INV)
    else:
        _, binary = cv2.threshold(gray, bg_gray + 30, 255, cv2.THRESH_BINARY)
    
    binary = cv2.morphologyEx(binary, cv2.MORPH_OPEN, np.ones((2,2), np.uint8))
    dist = cv2.distanceTransform(binary, cv2.DIST_L2, 5)
    
    # スケルトンまたは極大値からストローク幅を取得
    max_dist = np.max(dist)
    if max_dist <= 0: return 0.0, 0.0
    
    stroke_pixels = dist[dist > max_dist * 0.5]
    if len(stroke_pixels) == 0: return 0.0, 0.0
    
    avg_radius = np.mean(stroke_pixels)
    stroke_width = avg_radius * 2
    char_height = y2 - y1
    ratio = stroke_width / char_height if char_height > 0 else 0
    stroke_variation = np.std(stroke_pixels) / avg_radius if avg_radius > 0 else 0
    
    return ratio, stroke_variation

def estimate_font_size_v5(px_height, img_height, slide_height_pt):
    """v5: 画像の高さ比率からフォントサイズ(pt)を精密推定"""
    # 物理的なDPIに依存せず、スライドの論理高さ(pt)との比率で計算
    # 補正係数 0.70 は、bbox（上端から下端）に対する実際の「全角文字サイズ」の比率
    # 元JSONのデータ: height=60 in 1000scaleで font_size=32pt → 32/(768*60/1000)=0.69
    return (px_height / img_height) * slide_height_pt * 0.70

def refine_bbox_vertical(cv_image, bbox, bg_color):
    """テキストの垂直位置（Top/Bottom）を画素ベースで精密補正"""
    x1, y1, x2, y2 = bbox
    
    # 画像範囲チェック
    h, w = cv_image.shape[:2]
    x1, y1, x2, y2 = max(0, int(x1)), max(0, int(y1)), min(w, int(x2)), min(h, int(y2))
    
    if x2 <= x1 or y2 <= y1:
        return bbox
        
    roi = cv_image[y1:y2, x1:x2]
    
    # 背景色との差分
    diff = cv2.absdiff(roi, np.array(bg_color, dtype=np.uint8))
    gray_diff = cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)
    _, mask = cv2.threshold(gray_diff, 30, 255, cv2.THRESH_BINARY)
    
    # 行ごとの非ゼロピクセル数をカウント
    row_sums = cv2.reduce(mask, 1, cv2.REDUCE_SUM, dtype=cv2.CV_32S).flatten()
    
    # テキストが存在する行を探す
    y_indices = np.where(row_sums > 0)[0]
    
    if len(y_indices) > 0:
        true_top = y_indices[0]
        true_bottom = y_indices[-1]
        
        # マージン0: 画素境界をそのまま使用
        new_y1 = y1 + true_top
        new_y2 = y1 + true_bottom + 1
        
        return [x1, new_y1, x2, new_y2]
    
    return bbox

def get_font_for_text_v5(weight_ratio=0.0, stroke_var=0.0, color=(0, 0, 0)):
    """v5: 特徴量から日本語フォント候補を選択"""
    # 判定基準を実データに基づいて調整
    # デバッグログより: ratio=0.14～0.54 の範囲であった
    # タイトル(太字)はratio=0.14, 説明文はratio=0.40程度
    # この結果から、太字判定の閾値を上げる必要がある
    is_bold = weight_ratio > 0.30  # 0.12から引き上げ
    is_mincho = stroke_var > 0.5
    
    if is_mincho:
        font = "游明朝 Demibold" if is_bold else "游明朝"
    else:
        # ゴシック体: 太字ならHGPゴシックE、通常なら游ゴシック
        if is_bold:
            font = "HGPｺﾞｼｯｸE" 
        else:
            font = "游ゴシック"
            
    return font, is_bold

def get_background_color(image):
    """画像から背景色を検出"""
    w, h = image.size
    samples = [
        image.getpixel((10, 10)),
        image.getpixel((w-11, 10)),
        image.getpixel((10, h-11)),
        image.getpixel((w-11, h-11)),
        image.getpixel((w//2, 10)),
        image.getpixel((w//2, h-11)),
    ]
    samples = [s[:3] if len(s) > 3 else s for s in samples]
    return Counter(samples).most_common(1)[0][0]

def extract_text_blocks_native(page):
    """PyMuPDFでテキストブロックを抽出"""
    text_dict = page.get_text("dict")
    blocks = []
    
    for block in text_dict.get("blocks", []):
        if block.get("type") == 0:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if text:
                        bbox = span.get("bbox", [0, 0, 0, 0])
                        font_size = span.get("size", 12)
                        font_name = span.get("font", "")
                        color = span.get("color", 0)
                        
                        r = (color >> 16) & 0xFF
                        g = (color >> 8) & 0xFF
                        b = color & 0xFF
                        
                        blocks.append({
                            "text": text,
                            "bbox": bbox,
                            "font_size": font_size,
                            "font_name": font_name,
                            "color": (r, g, b),
                            "origin": span.get("origin", (bbox[0], bbox[1]))
                        })
    return blocks

def extract_text_with_ocr_word_level(page_image_path, page_width, page_height):
    """pytesseractでOCR - 単語レベルで抽出し、近接単語を結合"""
    import pytesseract
    from PIL import Image
    
    tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if os.path.exists(tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
    
    base_dir = os.path.dirname(os.path.abspath(__file__))
    # TesseractインストールフォルダからTESSDATA_PREFIXを設定（日本語パス回避）
    tesseract_tessdata = r"C:\Program Files\Tesseract-OCR\tessdata"
    if os.path.exists(tesseract_tessdata):
        os.environ['TESSDATA_PREFIX'] = tesseract_tessdata
        local_tessdata = tesseract_tessdata
    else:
        local_tessdata = os.path.join(base_dir, 'tessdata')
        os.environ['TESSDATA_PREFIX'] = local_tessdata
    print(f"  [DEBUG] Extracting text with OCR (word level)...")
    
    global OCR_SEGMENTS  # v35: グローバル変数を関数先頭で宣言
    OCR_SEGMENTS = []  # v35: セグメントリストを初期化
    
    # 画像読み込み (日本語パス対応)
    img_cv = read_image_cv2(page_image_path)
    if img_cv is None:
        print(f"Error: Failed to load image from {page_image_path}")
        return [], []
        
    # PIL形式に変換 (OCR用)
    img = Image.fromarray(cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB))
        
    img_width, img_height = img.size
    print(f"Using tessdata: {local_tessdata}")

    try:
        # OCRで詳細データを取得
        data = pytesseract.image_to_data(img, lang='jpn+eng', output_type=pytesseract.Output.DICT)
        
        # 単語単位で収集
        words = []
        for i in range(len(data['text'])):
            text = data['text'][i].strip()
            conf = data['conf'][i]
            
            # 信頼度-1は無視（ヘッダー行）
            if conf == -1:
                continue
                
            if text and conf > 50:  # 信頼度50%以上
                word_width = data['width'][i]
                word_height = data['height'][i]
                
                words.append({
                    'text': text,
                    'left': data['left'][i],
                    'top': data['top'][i],
                    'right': data['left'][i] + word_width,
                    'bottom': data['top'][i] + word_height,
                    'conf': conf,
                    'block_num': data['block_num'][i],
                    'line_num': data['line_num'][i],
                    'word_num': data['word_num'][i]
                })
        
        print(f"Found {len(words)} words with OCR")
        
        # グローバルに全ワードリストを保存（JSONブロック処理時に参照）
        global OCR_ALL_WORDS
        OCR_ALL_WORDS = words

        
        # ブロック/行単位でグループ化
        lines = defaultdict(list)
        for word in words:
            key = (word['block_num'], word['line_num'])
            lines[key].append(word)
        
        blocks = []
        text_bboxes_pixel = []
        
        for key, line_words in lines.items():
            if not line_words:
                continue
            
            # 単語を左から右にソート
            line_words.sort(key=lambda w: w['left'])
            
            # 行内で異常に大きい記号をフィルタリング
            # 条件：基準高さ（下位25%タイル）の2倍以上 AND 記号文字
            if len(line_words) > 1:
                heights = [w['bottom'] - w['top'] for w in line_words]
                sorted_heights = sorted(heights)
                # 下位25%タイルを基準に使用（異常値の影響を受けにくい）
                base_height = sorted_heights[max(0, len(heights) // 4)]
                
                symbol_chars = r'/\|-_=+()[]{}「」『』【】'
                filtered_words = []
                for w in line_words:
                    h = w['bottom'] - w['top']
                    # 基準の2倍以上 AND 単一記号文字 → スキップ
                    if h > base_height * 2 and len(w['text']) == 1 and w['text'] in symbol_chars:
                        print(f"  [DEBUG] Filtered out symbol: {w['text']} (h={h} > {base_height}*2)")
                        continue
                    filtered_words.append(w)
                line_words = filtered_words
            
            # 水平間隔が大きい場合、行を分割する
            split_lines = []
            current_segment = []
            
            for i, word in enumerate(line_words):
                if i == 0:
                    current_segment.append(word)
                else:
                    prev_right = line_words[i-1]['right']
                    curr_left = word['left']
                    gap = curr_left - prev_right
                    word_height = word['bottom'] - word['top']
                    
                    # 間隔が文字高さの2.5倍以上なら分割
                    if gap > word_height * 2.5:
                        if current_segment:
                            split_lines.append(current_segment)
                        current_segment = [word]
                    else:
                        current_segment.append(word)
            
            if current_segment:
                split_lines.append(current_segment)
            
            # 分割された各セグメントを個別のブロックとして処理
            for segment_idx, segment_words in split_lines:
                # 詳細情報を保存
                word_details = []
                combined_text = ""
                
                print(f"  [DEBUG] Processing segment {segment_idx} with {len(segment_words)} words")
                
                for i, word in enumerate(segment_words):
                    word_text = word['text']
                    print(f"    Word: '{word_text}' bbox={word['left']},{word['top']},{word['right']},{word['bottom']}")
                    
                    if i > 0:
                        # 前の単語との間隔チェック
                        prev_right = segment_words[i-1]['right']
                        curr_left = word['left']
                        gap = curr_left - prev_right
                        
                        # 隙間がある場合（文字高さの100%以上の場合のみスペース追加）
                        # 日本語テキストではスペースは基本的に不要
                        word_height = word['bottom'] - word['top']
                        if gap > word_height * 1.0:
                            print(f"      -> Adding space (gap={gap} > {word_height})")
                            word_details.append({
                                'text': ' ',
                                'bbox_pixel': None
                            })
                            combined_text += ' '
                    
                    word_text = word['text']
                    word_text = word_text.replace('`', '「').replace("'", '」')
                    
                    # 正規化（NFKC）：丸数字を数字に、全角英数を半角になど
                    word_text = unicodedata.normalize('NFKC', word_text)
                    
                    word_details.append({
                        'text': word_text,
                        'bbox_pixel': [word['left'], word['top'], word['right'], word['bottom']]
                    })
                    combined_text += word_text
                
                # print(f"    -> Combined text: '{combined_text}'")
                
                if len(combined_text.strip()) < 1:
                    # print("    -> Skipping empty block")
                    continue
                
                # バウンディングボックス計算
                x_coords = [d['bbox_pixel'][0] for d in word_details if d['bbox_pixel']]
                y_coords = [d['bbox_pixel'][1] for d in word_details if d['bbox_pixel']]
                r_coords = [d['bbox_pixel'][2] for d in word_details if d['bbox_pixel']]
                b_coords = [d['bbox_pixel'][3] for d in word_details if d['bbox_pixel']]
                
                if not x_coords:
                    continue
                    
                bbox = [min(x_coords), min(y_coords), max(r_coords), max(b_coords)]
                
                # フォントサイズ推定 (高さの80%)
                height = bbox[3] - bbox[1]
                font_size = height * 0.8
                
                # 色推定 (最初の単語の中心色)
                center_x = (bbox[0] + bbox[2]) // 2
                center_y = (bbox[1] + bbox[3]) // 2
                # Ensure coordinates are within image bounds
                center_x = max(0, min(center_x, img_cv.shape[1] - 1))
                center_y = max(0, min(center_y, img_cv.shape[0] - 1))
                color_bgr = img_cv[int(center_y), int(center_x)].tolist()
                color = [color_bgr[2], color_bgr[1], color_bgr[0]] # Convert BGR to RGB
                
                # 座標変換 (pixel -> point)
                # PDF座標系への変換
                pdf_bbox = [
                    bbox[0] / img_cv.shape[1] * page_width,
                    bbox[1] / img_cv.shape[0] * page_height,
                    bbox[2] / img_cv.shape[1] * page_width,
                    bbox[3] / img_cv.shape[0] * page_height
                ]
                
                blocks.append({
                    "text": combined_text,
                    "bbox": pdf_bbox,
                    "font_size": font_size,
                    "color": color,
                    "word_details": word_details
                })
                
                # テキスト領域リストに追加 (pixel座標)
                text_bboxes_pixel.append(bbox)
                
                # v35: セグメント情報をグローバルに保存
                OCR_SEGMENTS.append({
                    'text': combined_text,
                    'bbox_pixel': bbox,
                    'segment_words': segment_words
                })
        
        print(f"Grouped into {len(blocks)} text lines")
    except Exception as e:
        print(f"  [WARN] OCR processing error: {e}")
        # OCRエラーでも続行（補正データでのリカバリを試みる）

    # ★補正データ適用 (New Logic: Support    # 1. 補正データの読み込み
    corrections = load_ocr_corrections()
    page_key = f"page_{CURRENT_PAGE_NUM}"  # 動的にページ番号を設定
    
    if page_key in corrections:
        page_data = corrections[page_key]
        
        # 構造チェック: リストか辞書か
        if isinstance(page_data, list):
            # 互換性: 旧形式(リスト)の場合は既存ロジック(マージ)
            correction_list = page_data
            replace_all = False
        elif isinstance(page_data, dict):
            # 新形式
            correction_list = page_data.get("blocks", [])
            replace_all = page_data.get("replace_all", False)
        
        if replace_all:
            print(f"  [CORRECTION] REPLACING ALL text blocks for {page_key}")
            blocks = []     # 既存OCR全消去
            text_bboxes_pixel = [] # 重複チェック用もクリア
        else:
            print(f"  [CORRECTION] Merging corrections for {page_key}")

        for corr in correction_list:
            text = corr["text"]
            
            # v27: 新形式のcolors配列対応（マルチカラー）
            # 単一colorキーがある場合は従来通り、なければcolors配列を使用
            if "colors" in corr and isinstance(corr["colors"], list) and len(corr["colors"]) > 0:
                # 最初の色をデフォルト色として使用
                first_color = corr["colors"][0].get("rgb", [0, 0, 0])
                color = first_color
                multi_colors = corr["colors"]  # マルチカラー情報を保持
            else:
                color = corr.get("color", [0, 0, 0])
                multi_colors = None
            
            # JSONから座標を取得
            if "bbox_1000" in corr:
                nx, ny, nw, nh = corr["bbox_1000"]
                # 正規化座標(1000x1000) -> PDFポイント座標
                px = nx / 1000.0 * page_width
                py = ny / 1000.0 * page_height
                pw = nw / 1000.0 * page_width
                ph = nh / 1000.0 * page_height
                
                pdf_bbox = [px, py, px + pw, py + ph]
                
                # ピクセル座標（OCR検索のヒント用）
                json_px_bbox = [
                    int(nx / 1000.0 * img_cv.shape[1]),
                    int(ny / 1000.0 * img_cv.shape[0]),
                    int((nx + nw) / 1000.0 * img_cv.shape[1]),
                    int((ny + nh) / 1000.0 * img_cv.shape[0])
                ]
                json_bbox_pixel = [json_px_bbox[0], json_px_bbox[1], json_px_bbox[2], json_px_bbox[3]]
            else:
                # bbox_1000がない場合（万が一）
                pdf_bbox = [0, 0, 100, 20] # ダミー
                json_bbox_pixel = None
            
            # OCRから対応する座標を検索（ヒントを使用）
            print(f"  [v35 MATCH] Searching OCR for: '{text[:40]}...'")
            ocr_bbox = find_ocr_bbox_for_text(text, OCR_ALL_WORDS, json_bbox_pixel)
            
            final_bbox = pdf_bbox # デフォルトはJSON座標
            ocr_matched = False
            matched_segment_text = None  # v35: マッチしたセグメントのテキスト
            
            if ocr_bbox:
                # OCRで見つかった場合は、位置をOCRベースに更新
                px_x1, px_y1, px_x2, px_y2 = ocr_bbox
                
                # v35: マッチしたセグメントを探す（ocr_bboxとの重なりが最大のものを選択）
                best_segment = None
                max_overlap_area = 0
                
                # ocr_bbox: [x1, y1, x2, y2]
                ocr_area = (px_x2 - px_x1) * (px_y2 - px_y1)
                
                for seg in OCR_SEGMENTS:
                    seg_bbox = seg['bbox_pixel'] # [sx1, sy1, sx2, sy2]
                    
                    # 重なり計算 (Intersection)
                    ix1 = max(px_x1, seg_bbox[0])
                    iy1 = max(px_y1, seg_bbox[1])
                    ix2 = min(px_x2, seg_bbox[2])
                    iy2 = min(px_y2, seg_bbox[3])
                    
                    if ix2 > ix1 and iy2 > iy1:
                        intersection = (ix2 - ix1) * (iy2 - iy1)
                        if intersection > max_overlap_area:
                            max_overlap_area = intersection
                            best_segment = seg
                
                # 重なりが小さすぎる場合はマッチなしとみなす（ノイズ対策）
                if ocr_area > 0 and (max_overlap_area / ocr_area) < 0.3:
                    best_segment = None
                
                if best_segment:
                    matched_segment_text = best_segment['text']
                    segment_bbox = best_segment['bbox_pixel']
                    
                    # v35: 補正判定 - JSONテキストとセグメントテキストの先頭/末尾を比較
                    json_first = text[0] if text else ''
                    json_last = text[-1] if text else ''
                    ocr_first = matched_segment_text[0] if matched_segment_text else ''
                    ocr_last = matched_segment_text[-1] if matched_segment_text else ''
                    
                    first_match = (json_first == ocr_first)
                    last_match = (json_last == ocr_last)
                    
                    # v36: Anchor to JSON Top-Left, Expand Bottom-Right with max
                    # JSONの左上を使用（レイアウト意図を保持）
                    
                    seg_x1, seg_y1, seg_x2, seg_y2 = segment_bbox
                    
                    if json_bbox_pixel:
                        json_x1, json_y1, json_x2, json_y2 = json_bbox_pixel
                    else:
                        json_x1, json_y1, json_x2, json_y2 = segment_bbox
                    
                    # v36: 左上はJSON、右下はmax(JSON, OCR)
                    px_x1 = json_x1
                    px_y1 = json_y1
                    px_x2 = max(json_x2, seg_x2)
                    px_y2 = max(json_y2, seg_y2)
                    
                    print(f"  [v36] Text: '{text[:30]}...'")
                    print(f"  [v36] Seg(OCR): {seg_x1},{seg_y1}-{seg_x2},{seg_y2}")
                    print(f"  [v36] JSON:     {json_x1},{json_y1}-{json_x2},{json_y2}")
                    print(f"  [v36] Merged:   {px_x1},{px_y1}-{px_x2},{px_y2} (Anchor JSON Top-Left)")

                    # 2. 文字不足補正 (OCRセグメント自体がJSONより短い場合)
                    if not (first_match and last_match):
                        json_len = len(text)
                        ocr_len = len(matched_segment_text)
                        
                        if ocr_len > 0 and json_len > ocr_len:
                            missing_chars = json_len - ocr_len
                            # 現在の幅から1文字幅を推定
                            current_width = px_x2 - px_x1
                            char_width = current_width / ocr_len if ocr_len > 0 else 20
                            additional_width = char_width * missing_chars
                            
                            if first_match and not last_match:
                                # 末尾漏れ → 右にさらに拡張
                                px_x2 += additional_width
                                print(f"  [v35] Missing chars at end ({missing_chars}). Extending Right +{additional_width:.1f}px")
                            elif not first_match and last_match:
                                # 先頭漏れ → 左には拡張せず、右に拡張するか？あるいは左を少し広げる？
                                # 左に広げると「左ズレ」問題が再発するおそれがあるが、文字が切れるよりマシか？
                                # 今回のユーザー指摘は「左にずれて要素を消した」なので、左拡張は極力避けるべき。
                                # JSONに含まれていない文字があるなら、JSONが間違っているか、OCRがノイズ。
                                # しかしJSONは正解テキストを持っている。
                                # ここでは「左への拡張」は最大でも文字幅程度に留めるか、しない。
                                pass # 先頭漏れは稀であり、左ズレリスクが高いため拡張しない（JSONのx1を信じる）
                            else:
                                # 両方不一致
                                pass
                
                # 最終確認ログ
                print(f"  [v35] Final Pixel BBox: {int(px_x1)},{int(px_y1)},{int(px_x2)},{int(px_y2)}")
                
                # PDF座標系への変換
                final_bbox = [
                    px_x1 / img_cv.shape[1] * page_width,
                    px_y1 / img_cv.shape[0] * page_height,
                    px_x2 / img_cv.shape[1] * page_width,
                    px_y2 / img_cv.shape[0] * page_height
                ]
                
                # ピクセル座標をtext_bboxes_pixelに追加
                pixel_bbox = [px_x1, px_y1, px_x2, px_y2]
                text_bboxes_pixel.append(pixel_bbox)
                print(f"  [CORRECTION] Match found. Using Merged bbox.")
                ocr_matched = True
            else:
                # OCRで見つからない場合はJSON座標を使用（フォールバック）
                if json_bbox_pixel:
                    text_bboxes_pixel.append(json_bbox_pixel)
                print(f"  [CORRECTION] '{text[:30]}...' -> No OCR match. Using JSON bbox (fallback).")
            
            # v27: フォントサイズ（font_size_ptキーを優先）
            if "font_size_pt" in corr:
                font_size = corr["font_size_pt"]
            elif "font_size_fixed" in corr:
                font_size = corr["font_size_fixed"]
            else:
                # フォールバック: 高さの80%
                font_size = (final_bbox[3] - final_bbox[1]) * 0.8
            
            # v36: フォントサイズに1.2倍の補正係数を適用
            font_size = font_size * 1.2
            print(f"  [v36] Font size: {font_size/1.2:.1f}pt -> {font_size:.1f}pt (1.2x corrected)")
            
            # v27: フォントファミリー
            font_family = corr.get("font_family", "Noto Sans JP")
            
            # v27: 太字判定（JSONから取得）
            is_bold = corr.get("is_bold", False)
            
            # ブロック追加
            # 色情報の取得 (v39 fix)
            json_color = [0, 0, 0] # デフォルト黒
            if "colors" in corr and corr["colors"]:
                # 最初の色を使用（あるいは最も範囲が広い色を選ぶロジックでも良いが、通常は単色が支配的）
                try:
                    first_color = corr["colors"][0]
                    if "rgb" in first_color:
                        json_color = first_color["rgb"]
                    print(f"  [v39 DEBUG] Extracted JSON color: {json_color}")
                except Exception as e:
                    print(f"  [ERROR] Failed to extract color from JSON: {e}")
            else:
                print(f"  [v39 DEBUG] No 'colors' array found in JSON, using default black or 'color' key.")
                json_color = corr.get("color", [0, 0, 0]) # Fallback to old 'color' key
            
            block_data = {
                "text": text,
                "bbox": final_bbox,
                "font_size": font_size,
                "color": json_color, # v39: JSONから取得した色を優先
                "word_details": [],
                "font_family": font_family,
                "is_bold": is_bold,
                "ocr_matched": ocr_matched,
                "json_font_size": corr.get("font_size_pt", font_size)  # 後のフォントサイズ補正用
            }
            
            # マルチカラー情報があれば追加
            if multi_colors:
                block_data["multi_colors"] = multi_colors
            
            # オプショナルキーの継承（互換性）
            if "background_color" in corr: block_data["background_color"] = corr["background_color"]
            
            blocks.append(block_data)

        return blocks, text_bboxes_pixel

def detect_text_color(img, bbox, bg_color):
    """テキスト領域の支配的な色を検出（最頻色・アンチエイリアスを除外）"""
    # bbox = [x1, y1, x2, y2] (左, 上, 右, 下)
    x1, y1, x2, y2 = [int(v) for v in bbox]
    h, w = img.shape[:2]  # h=高さ(rows), w=幅(cols)
    
    # 境界チェック（x: 0〜w, y: 0〜h）
    x1, x2 = max(0, x1), min(w, x2)
    y1, y2 = max(0, y1), min(h, y2)
    
    if x2 <= x1 or y2 <= y1:
        return (255, 255, 255)
    
    # OpenCVは[row, col] = [y, x]の順序
    region = img[y1:y2, x1:x2]
    
    # 背景色と異なるピクセルを抽出
    bg_array = np.array(bg_color)
    diff = np.abs(region.astype(np.float32) - bg_array)
    diff_sum = np.sum(diff, axis=2)
    
    # 背景から十分離れたピクセル（diff_sum > 80でアンチエイリアスを除外）
    mask = diff_sum > 80
    
    if np.any(mask):
        text_pixels = region[mask]
        
        # 最も暗い色（最小輝度）を検出（テキストは通常最も暗い）
        # 輝度 = 0.299*R + 0.587*G + 0.114*B
        luminances = 0.299 * text_pixels[:, 2] + 0.587 * text_pixels[:, 1] + 0.114 * text_pixels[:, 0]
        
        # 最も暗い10%のピクセルの平均色を使用
        dark_threshold = np.percentile(luminances, 10)
        dark_mask = luminances <= dark_threshold
        dark_pixels = text_pixels[dark_mask]
        
        if len(dark_pixels) > 0:
            avg_color = np.mean(dark_pixels, axis=0)
            # BGR→RGB変換
            r, g, b = [max(0, min(255, int(c))) for c in avg_color[::-1]]
            return (r, g, b)
        
        # フォールバック: 全体の平均
        avg_color = np.mean(text_pixels, axis=0)
        r, g, b = [max(0, min(255, int(c))) for c in avg_color[::-1]]
        return (r, g, b)
    
    return (255, 255, 255)

def is_text_region(elem_bbox, text_bboxes, threshold=0.1):
    """
    グラフィック要素がテキスト領域と重なっているか判定
    elem_bbox: (x, y, w, h)
    text_bboxes: [[x1, y1, x2, y2], ...] (pixel coordinates)
    """
    elem_x, elem_y, elem_w, elem_h = elem_bbox
    elem_area = elem_w * elem_h
    if elem_area <= 0:
        return False
        
    elem_right = elem_x + elem_w
    elem_bottom = elem_y + elem_h
    
    for text_bbox in text_bboxes:
        # text_bbox is now [x1, y1, x2, y2]
        if isinstance(text_bbox, dict):
            # 古い形式互換（念のため）
            text_left = text_bbox['left']
            text_top = text_bbox['top']
            text_right = text_bbox['right']
            text_bottom = text_bbox['bottom']
        else:
            # 新しいリスト形式
            text_left = text_bbox[0]
            text_top = text_bbox[1]
            text_right = text_bbox[2]
            text_bottom = text_bbox[3]
            
        # 重なり判定
        overlap_left = max(elem_x, text_left)
        overlap_top = max(elem_y, text_top)
        overlap_right = min(elem_right, text_right)
        overlap_bottom = min(elem_bottom, text_bottom)
        
        if overlap_right > overlap_left and overlap_bottom > overlap_top:
            overlap_area = (overlap_right - overlap_left) * (overlap_bottom - overlap_top)
            # 要素の面積に対する重なり割合
            if overlap_area / elem_area > threshold:
                # 重なっているが、要素の面積がテキストBBoxの面積よりも遥かに大きい（1.5倍以上）なら
                # それは「背景ボックス」である可能性が高いので、テキストとみなさず（Falseを返し）救済する
                text_width = text_right - text_left
                text_height = text_bottom - text_top
                text_area = text_width * text_height
                if text_area > 0 and elem_area > text_area * 1.1:
                    return False
                
                return True
                
    return False

def recover_text_from_crop(img_crop):
    """画像クロップに対してOCRを実行し、テキストが含まれているか確認"""
    try:
        # 画像が小さすぎる場合はスキップ
        if img_crop.shape[0] < 10 or img_crop.shape[1] < 10:
            return None
            
        # パディング（余白）を追加 - 大きめに取る
        padding = 40
        h, w = img_crop.shape[:2]
        padded_img = np.ones((h + 2 * padding, w + 2 * padding, 3), dtype=np.uint8) * 255
        padded_img[padding:padding+h, padding:padding+w] = img_crop
        
        # ★アップスケール（拡大）: 小さい文字の認識率向上
        scale = 3
        padded_img = cv2.resize(padded_img, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
        
        pil_img = Image.fromarray(cv2.cvtColor(padded_img, cv2.COLOR_BGR2RGB))
        
        # 1. ブロックとしてOCR (--psm 6)
        text = pytesseract.image_to_string(pil_img, lang='jpn', config='--psm 6').strip()
        if len(text) > 0:
            print(f"  [RECOVERY] Recovered text (psm 6): '{text}'")
            return text

        # 2. 1文字としてOCR (--psm 10) - 1文字の場合
        text = pytesseract.image_to_string(pil_img, lang='jpn', config='--psm 10').strip()
        if len(text) > 0:
            print(f"  [RECOVERY] Recovered text (psm 10): '{text}'")
            return text
            
        # デバッグ用：失敗した画像を保存（特定の座標周辺のみ）
        # pixel bbox info required here? No, just save.
        # But we don't have coordinates here easily. Just save random name.
        # Limit to small images to avoid disk spam
        if h < 100 and w < 100:
             cv2.imwrite(f"debug_fail_recovery_{np.random.randint(10000)}.png", padded_img)
            
    except Exception as e:
        pass
    return None

def detect_graphic_elements(page_image_path, output_dir, text_bboxes_pixel, text_bboxes_with_colors=None):
    """OpenCVで輪郭検出 - テキスト領域と一定以上重なる要素を除外
    v8: ネスト検出 - 素材内のテキストボックス重複時、内部で再度素材検出を行い、
        検出された部分だけをマスク。境界外の色でマスク塗りつぶし。
    v39: text_bboxes_with_colors = [(bbox, color_rgb), ...] でテキスト色情報を受け取る
    """
    img = read_image_cv2(page_image_path)
    if img is None:
        print(f"Failed to read image: {page_image_path}")
        return []
    
    print(f"Image loaded: {img.shape}")
    
    # 背景色を取得
    bg_color = img[10, 10].copy()
    print(f"Background color (BGR): {bg_color}")
    
    # 背景との差分でマスク作成
    diff = cv2.absdiff(img, np.full_like(img, bg_color))
    diff_gray = cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)
    _, element_mask = cv2.threshold(diff_gray, 25, 255, cv2.THRESH_BINARY)
    
    # ノイズ除去
    kernel = np.ones((3, 3), np.uint8)
    element_mask = cv2.morphologyEx(element_mask, cv2.MORPH_CLOSE, kernel)
    element_mask = cv2.morphologyEx(element_mask, cv2.MORPH_OPEN, kernel)
    
    # 輪郭検出
    contours, _ = cv2.findContours(element_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    print(f"Found {len(contours)} total contours")
    
    elements = []
    recovered_texts = []
    min_area = 500
    text_overlap_count = 0
    nested_detect_count = 0
    
    def find_contained_text_boxes(elem_bbox, text_bboxes):
        """グラフィック要素が完全に含むテキストボックスを検索"""
        ex, ey, ew, eh = elem_bbox
        contained = []
        for tb in text_bboxes:
            tx1, ty1, tx2, ty2 = tb
            text_w = tx2 - tx1
            text_h = ty2 - ty1
            text_area = text_w * text_h
            if text_area <= 0:
                continue
            
            overlap_x1 = max(ex, tx1)
            overlap_y1 = max(ey, ty1)
            overlap_x2 = min(ex + ew, tx2)
            overlap_y2 = min(ey + eh, ty2)
            
            if overlap_x2 > overlap_x1 and overlap_y2 > overlap_y1:
                overlap_area = (overlap_x2 - overlap_x1) * (overlap_y2 - overlap_y1)
                if overlap_area / text_area > 0.5:  # v16: 80%から50%に下げ
                    contained.append(tb)
        return contained

    def get_boundary_color(img_crop, margin=3):
        """素材の境界外の色を取得（上端の少し外）"""
        h, w = img_crop.shape[:2]
        # 上端のmarginピクセル外の色をサンプリング
        if margin < h:
            # 上端のライン
            top_colors = img_crop[0:min(margin, h), :, :]
            avg_color = np.mean(top_colors.reshape(-1, 3), axis=0)
            return avg_color.astype(np.uint8)
        return np.array([255, 255, 255], dtype=np.uint8)

    def detect_sub_elements_in_crop(crop_img, crop_bg_color):
        """クロップ画像内で素材を再検出"""
        diff = cv2.absdiff(crop_img, np.full_like(crop_img, crop_bg_color))
        diff_gray = cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)
        _, sub_mask = cv2.threshold(diff_gray, 75, 255, cv2.THRESH_BINARY)
        
        kernel = np.ones((2, 2), np.uint8)
        sub_mask = cv2.morphologyEx(sub_mask, cv2.MORPH_CLOSE, kernel)
        sub_mask = cv2.morphologyEx(sub_mask, cv2.MORPH_OPEN, kernel)
        
        # v10: RETR_LIST で箱の中のテキスト要素も検出（RETR_EXTERNALは外側の枠しか取らない）
        sub_contours, _ = cv2.findContours(sub_mask, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
        
        sub_elements = []
        for contour in sub_contours:
            area = cv2.contourArea(contour)
            if area < 30:  # サブ要素の最小面積
                continue
            sx, sy, sw, sh = cv2.boundingRect(contour)
            sub_elements.append({
                'bbox': (sx, sy, sw, sh),
                'area': area
            })
        return sub_elements, sub_mask
    
    for i, contour in enumerate(contours):
        area = cv2.contourArea(contour)
        if area < min_area:
            continue
        
        x, y, w, h = cv2.boundingRect(contour)
        
        # 要素を切り出し
        element = img[y:y+h, x:x+w].copy()
        element_mask_crop = element_mask[y:y+h, x:x+w].copy()
        
        # リカバリーOCR
        recovered_text = recover_text_from_crop(element)
        
        if 2200 < x < 2300 and 750 < y < 850:
             print(f"  [DEBUG_ZO] Contour at ({x},{y},{w},{h}), Area={area}")
             print(f"  [DEBUG_ZO] Recovered text: '{recovered_text}'")

        if recovered_text:
            recovered_texts.append({
                "text": recovered_text,
                "bbox": [x, y, x+w, y+h],
                "word_details": [] 
            })
            continue
        
        # ★v8: テキストボックスが80%以上含まれている場合、ネスト検出
        contained_text_boxes = find_contained_text_boxes((x, y, w, h), text_bboxes_pixel)
        
        if contained_text_boxes:
            print(f"  [NESTED] Graphic at ({x},{y},{w},{h}) contains {len(contained_text_boxes)} text boxes")
            
            # 素材内の背景色を推定（中央付近の色をサンプリング - 枠線ではなく内部の色を取得）
            # 中央から少し内側の複数点をサンプリングして最頻値を使う
            sample_points = [
                (y + h // 4, x + w // 4),      # 左上1/4
                (y + h // 4, x + 3 * w // 4),  # 右上1/4
                (y + h // 2, x + w // 2),      # 中央
                (y + 3 * h // 4, x + w // 4),  # 左下1/4
                (y + 3 * h // 4, x + 3 * w // 4),  # 右下1/4
            ]
            colors = []
            for sy, sx in sample_points:
                if 0 <= sy < img.shape[0] and 0 <= sx < img.shape[1]:
                    colors.append(tuple(img[sy, sx]))
            
            if colors:
                # 最頻値を使用
                from collections import Counter
                inner_bg_color = np.array(Counter(colors).most_common(1)[0][0], dtype=np.uint8)
            else:
                inner_bg_color = bg_color
            
            print(f"    [NESTED] Inner bg color sampled: {inner_bg_color}")
            
            # v25: サブ要素検出はparent_fill_color計算後に行う（後で呼び出し）
            # 素材内で再度素材検出 - 一時的にinner_bg_colorを使用してparent_fill_colorを計算後に再検出
            
            # v13: 各サブ要素ごとに周囲8点から背景色をサンプリング（塗りつぶし前に収集）
            crop_h, crop_w = element.shape[:2]
            parent_area = crop_h * crop_w
            
            # v20: 親要素の内側境界から10px内側を32点周回サンプリング
            def sample_interior_border(img, inset=10, points_per_side=8):
                """親要素の内側境界から指定ピクセル内側を周回サンプリング"""
                h, w = img.shape[:2]
                colors = []
                
                # 内側の境界座標を計算
                x1 = inset
                y1 = inset
                x2 = w - inset
                y2 = h - inset
                
                # 境界が有効かチェック
                if x2 <= x1 or y2 <= y1:
                    # 要素が小さすぎる場合は中央の色を返す
                    cy, cx = h // 2, w // 2
                    if 0 <= cy < h and 0 <= cx < w:
                        return [tuple(img[cy, cx])]
                    return []
                
                # 上辺: points_per_side点
                for i in range(points_per_side):
                    px = x1 + int((x2 - x1) * (i + 1) / (points_per_side + 1))
                    py = y1
                    if 0 <= py < h and 0 <= px < w:
                        colors.append(tuple(img[py, px]))
                
                # 下辺: points_per_side点
                for i in range(points_per_side):
                    px = x1 + int((x2 - x1) * (i + 1) / (points_per_side + 1))
                    py = y2
                    if 0 <= py < h and 0 <= px < w:
                        colors.append(tuple(img[py, px]))
                
                # 左辺: points_per_side点
                for i in range(points_per_side):
                    px = x1
                    py = y1 + int((y2 - y1) * (i + 1) / (points_per_side + 1))
                    if 0 <= py < h and 0 <= px < w:
                        colors.append(tuple(img[py, px]))
                
                # 右辺: points_per_side点
                for i in range(points_per_side):
                    px = x2
                    py = y1 + int((y2 - y1) * (i + 1) / (points_per_side + 1))
                    if 0 <= py < h and 0 <= px < w:
                        colors.append(tuple(img[py, px]))
                
                return colors
            
            def cluster_colors_and_get_dominant(colors, threshold=40):
                """v18: 色をクラスタリングして最も多いクラスタの代表色を返す"""
                if not colors:
                    return None
                
                # クラスタリング: 類似色（距離threshold以内）をグループ化
                clusters = []  # [[color1, color2, ...], ...]
                
                for color in colors:
                    found_cluster = False
                    for cluster in clusters:
                        # クラスタの代表色（最初の色）との距離を計算
                        representative = cluster[0]
                        dist = ((color[0] - representative[0])**2 + 
                                (color[1] - representative[1])**2 + 
                                (color[2] - representative[2])**2) ** 0.5
                        if dist < threshold:
                            cluster.append(color)
                            found_cluster = True
                            break
                    
                    if not found_cluster:
                        # 新しいクラスタを作成
                        clusters.append([color])
                
                # 最も要素数が多いクラスタを選択
                largest_cluster = max(clusters, key=len)
                
                # そのクラスタの平均色を計算
                avg_b = int(np.mean([c[0] for c in largest_cluster]))
                avg_g = int(np.mean([c[1] for c in largest_cluster]))
                avg_r = int(np.mean([c[2] for c in largest_cluster]))
                
                return np.array([avg_b, avg_g, avg_r], dtype=np.uint8)
            
            # v22: 親要素の内側境界から20px内側を32点周回サンプリング
            parent_interior_colors = sample_interior_border(element, inset=25, points_per_side=8)
            if len(parent_interior_colors) >= 3:
                parent_fill_color = cluster_colors_and_get_dominant(parent_interior_colors, threshold=2)
                if parent_fill_color is None:
                    parent_fill_color = inner_bg_color
            else:
                parent_fill_color = inner_bg_color
            
            print(f"    [NESTED] v25: Parent interior sampled {len(parent_interior_colors)} pts -> fill color BGR={parent_fill_color}")
            
            # v25: parent_fill_colorを使用してサブ要素を検出
            sub_elements, sub_mask = detect_sub_elements_in_crop(element, parent_fill_color)
            print(f"    [NESTED] Found {len(sub_elements)} sub-elements inside (using parent_fill_color)")
            
            # Step 1: サブ要素を収集（塗りつぶし色は親から取得済み）
            valid_sub_elements = []
            
            for sub_elem in sub_elements:
                sx, sy, sw, sh = sub_elem['bbox']
                sub_area = sw * sh
                
                # 親要素の85%以上を占めるサブ要素は枠線とみなして除外
                if sub_area > parent_area * 0.85:
                    print(f"    [NESTED] Skipped large sub-element (likely frame): local({sx},{sy},{sw},{sh}), area={sub_area}/{parent_area}")
                    continue
                
                # v15: サブ要素がテキストボックスを含むかチェック
                # サブ要素のグローバル座標を計算
                global_sx = x + sx
                global_sy = y + sy
                sub_contained_text = find_contained_text_boxes((global_sx, global_sy, sw, sh), text_bboxes_pixel)
                
                if sub_contained_text:
                    # v20: サブ親要素も内側境界サンプリングで色を決定
                    print(f"    [NESTED-L2] Sub-element at local({sx},{sy},{sw},{sh}) contains {len(sub_contained_text)} text boxes -> recursive detection")
                    
                    # サブ要素をクロップ
                    sub_crop = element[sy:sy+sh, sx:sx+sw].copy()
                    sub_crop_h, sub_crop_w = sub_crop.shape[:2]
                    sub_parent_area = sub_crop_h * sub_crop_w
                    
                    # v22: サブ親要素の内側境界から20pxサンプリング
                    sub_interior_colors = sample_interior_border(sub_crop, inset=20, points_per_side=8)
                    if len(sub_interior_colors) >= 3:
                        sub_fill_color = cluster_colors_and_get_dominant(sub_interior_colors, threshold=2)
                        if sub_fill_color is None:
                            sub_fill_color = parent_fill_color
                    else:
                        sub_fill_color = parent_fill_color
                    
                    print(f"      [NESTED-L2] Sub-parent interior sampled {len(sub_interior_colors)} pts -> fill color BGR={sub_fill_color}")
                    
                    # サブ要素内でさらにサブサブ要素を検出
                    sub_sub_elements, _ = detect_sub_elements_in_crop(sub_crop, sub_fill_color)
                    print(f"      [NESTED-L2] Found {len(sub_sub_elements)} sub-sub-elements")
                    
                    # v42: サブサブ要素をオーバーレイに追加（巨大なものは除外だが塗りつぶしはする）
                    for ss_elem in sub_sub_elements:
                        ssx, ssy, ssw, ssh = ss_elem['bbox']
                        ss_area = ssw * ssh
                        
                        # サブ要素の85%以上は枠として塗りつぶしもスキップ
                        if ss_area > sub_parent_area * 0.85:
                            print(f"      [v42] Skipped sub-sub (frame): {ss_area} > 85% of {sub_parent_area}")
                            continue
                        
                        # v42: サブ親の50%以上の巨大なサブサブ要素は塗りつぶしもオーバーレイもスキップ
                        if ss_area >= sub_parent_area * 0.5:
                            print(f"      [v42] Skipped large sub-sub: {ss_area} >= 50% of {sub_parent_area}")
                            continue
                        
                        # サブサブ要素を塗りつぶし（サブクロップ内）- サブ親の色を使用
                        fill_margin = 5
                        ss_mask_sx = max(0, ssx - fill_margin)
                        ss_mask_sy = max(0, ssy - fill_margin)
                        ss_mask_ex = min(sub_crop_w, ssx + ssw + fill_margin)
                        ss_mask_ey = min(sub_crop_h, ssy + ssh + fill_margin)
                        
                        sub_crop[ss_mask_sy:ss_mask_ey, ss_mask_sx:ss_mask_ex] = sub_fill_color
                        print(f"      [NESTED-L2] Filled sub-sub at ({ssx},{ssy},{ssw},{ssh}) with BGR={sub_fill_color}")
                        
                        # v42: 小さいサブサブ要素をオーバーレイとして追加
                        # サブサブ要素のグローバル座標を計算
                        global_ssx = x + sx + ssx
                        global_ssy = y + sy + ssy
                        
                        # v42 fix: サブサブ要素もテキスト重複チェックを適用（通常サブ要素と同様）
                        ss_text_overlap = is_text_region((global_ssx, global_ssy, ssw, ssh), text_bboxes_pixel, threshold=0.2)
                        
                        if ss_text_overlap:
                            # テキスト重複時は色差による免除判定
                            def get_overlapping_text_colors_ss(sub_bbox, text_bboxes_with_colors):
                                if not text_bboxes_with_colors:
                                    return []
                                ssx, ssy, ssw, ssh = sub_bbox
                                sub_right, sub_bottom = ssx + ssw, ssy + ssh
                                text_colors = []
                                for bbox, color_rgb in text_bboxes_with_colors:
                                    tx1, ty1, tx2, ty2 = bbox
                                    if max(ssx, tx1) < min(sub_right, tx2) and max(ssy, ty1) < min(sub_bottom, ty2):
                                        color_bgr = (color_rgb[2], color_rgb[1], color_rgb[0]) if color_rgb else (0, 0, 0)
                                        text_colors.append(color_bgr)
                                return text_colors
                            
                            ss_overlapping_colors = get_overlapping_text_colors_ss(
                                (global_ssx, global_ssy, ssw, ssh), text_bboxes_with_colors
                            )
                            
                            ss_exempt = False
                            for txt_color in ss_overlapping_colors:
                                ss_color_diff = abs(int(txt_color[0]) - int(sub_fill_color[0])) + \
                                                abs(int(txt_color[1]) - int(sub_fill_color[1])) + \
                                                abs(int(txt_color[2]) - int(sub_fill_color[2]))
                                if ss_color_diff <= 70:
                                    ss_exempt = True
                                    print(f"      [v42] Sub-sub text overlap: color diff={ss_color_diff} <= 70 -> EXEMPT")
                                    break
                            
                            if not ss_exempt:
                                print(f"      [v42] Sub-sub at ({global_ssx},{global_ssy},{ssw},{ssh}): text overlap, no color exempt -> skip overlay")
                                continue
                            else:
                                print(f"      [v42] Sub-sub at ({global_ssx},{global_ssy},{ssw},{ssh}): text overlap but COLOR EXEMPT")
                        else:
                            print(f"      [v42] Sub-sub at ({global_ssx},{global_ssy},{ssw},{ssh}): no text overlap -> add to overlay")
                        
                        # margin=7で拡張して切り出し
                        overlay_margin = 7
                        ss_crop_sx = max(0, global_ssx - overlay_margin)
                        ss_crop_sy = max(0, global_ssy - overlay_margin)
                        ss_crop_ex = min(img.shape[1], global_ssx + ssw + overlay_margin)
                        ss_crop_ey = min(img.shape[0], global_ssy + ssh + overlay_margin)
                        
                        ss_crop_img = img[ss_crop_sy:ss_crop_ey, ss_crop_sx:ss_crop_ex].copy()
                        ss_expanded_w = ss_crop_ex - ss_crop_sx
                        ss_expanded_h = ss_crop_ey - ss_crop_sy
                        
                        # RGBA変換
                        ss_rgba = cv2.cvtColor(ss_crop_img, cv2.COLOR_BGR2BGRA)
                        ss_rgba[:, :, 3] = 255
                        
                        ss_path = os.path.join(output_dir, f"subsub_element_{i}_{sx}_{sy}_{ssx}_{ssy}.png")
                        cv2.imwrite(ss_path, ss_rgba)
                        
                        elements.append({
                            "path": ss_path,
                            "bbox": (ss_crop_sx, ss_crop_sy, ss_expanded_w, ss_expanded_h),
                            "area": ss_expanded_w * ss_expanded_h,
                            "is_sub_overlay": True
                        })
                        print(f"      [v42] Added sub-sub overlay: {ss_path}")
                    
                    # サブ要素のクロップを親要素に書き戻す
                    element[sy:sy+sh, sx:sx+sw] = sub_crop
                    element_mask_crop[sy:sy+sh, sx:sx+sw] = 255
                    print(f"    [NESTED-L2] Wrote back processed sub-element to parent")
                    # v42: サブ親自体はオーバーレイに追加しない（サブサブが追加されるので）
                    continue
                
                # 通常のサブ要素（テキストを含まない）: valid_sub_elementsに追加
                valid_sub_elements.append(sub_elem)
                print(f"    [NESTED] Sub-element at local({sx},{sy},{sw},{sh}): will use parent fill color BGR={parent_fill_color}")
            
            # Step 2: 親の色で塗りつぶし実行
            masked_count = 0
            for sub_elem in valid_sub_elements:
                sx, sy, sw, sh = sub_elem['bbox']
                
                # マスク範囲を5ピクセル拡張（輪郭残り防止）
                margin = 5
                mask_sx = max(0, sx - margin)
                mask_sy = max(0, sy - margin)
                mask_ex = min(element.shape[1], sx + sw + margin)
                mask_ey = min(element.shape[0], sy + sh + margin)
                
                # v21: 親の塗りつぶし色で塗りつぶす
                element[mask_sy:mask_ey, mask_sx:mask_ex] = parent_fill_color
                # v14修正: 透明(0)ではなく不透明(255)にして塗りつぶした色を表示
                element_mask_crop[mask_sy:mask_ey, mask_sx:mask_ex] = 255
                masked_count += 1
                print(f"    [FILL] Applied BGR={parent_fill_color} to region ({mask_sx},{mask_sy})-({mask_ex},{mask_ey})")
            
            print(f"    [NESTED] Total masked: {masked_count}/{len(sub_elements)} sub-elements")
            
            # v37: サブ要素を別途オーバーレイグラフィックとして収集
            # 塗りつぶし前のオリジナル画像からサブ要素を切り出して別途保存
            def sample_8_points_outside_within_parent(full_img, global_x, global_y, sw, sh, parent_x, parent_y, parent_w, parent_h, offset=10):
                """サブ要素の10px外側8点をサンプリング（ただし親要素内に限定）"""
                h_img, w_img = full_img.shape[:2]
                
                # 親要素の境界
                parent_x2 = parent_x + parent_w
                parent_y2 = parent_y + parent_h
                
                points = [
                    # 上辺3点 (左, 中, 右)
                    (global_y - offset, global_x + sw // 4),
                    (global_y - offset, global_x + sw // 2),
                    (global_y - offset, global_x + 3 * sw // 4),
                    # 下辺3点
                    (global_y + sh + offset, global_x + sw // 4),
                    (global_y + sh + offset, global_x + sw // 2),
                    (global_y + sh + offset, global_x + 3 * sw // 4),
                    # 左辺1点
                    (global_y + sh // 2, global_x - offset),
                    # 右辺1点
                    (global_y + sh // 2, global_x + sw + offset),
                ]
                colors = []
                valid_count = 0
                for py, px in points:
                    # 画像境界内かつ親要素境界内の点のみサンプリング
                    if 0 <= py < h_img and 0 <= px < w_img:
                        if parent_x <= px < parent_x2 and parent_y <= py < parent_y2:
                            colors.append(full_img[py, px])
                            valid_count += 1
                return colors, valid_count
            
            def compute_color_diff(colors, base_color):
                """8点の色と基本色の平均色差を計算"""
                if not colors:
                    return 0
                diffs = []
                for c in colors:
                    diff = abs(int(c[0]) - int(base_color[0])) + \
                           abs(int(c[1]) - int(base_color[1])) + \
                           abs(int(c[2]) - int(base_color[2]))
                    diffs.append(diff)
                return np.mean(diffs)
            
            for sub_elem in valid_sub_elements:
                sx, sy, sw, sh = sub_elem['bbox']
                
                # サブ要素のグローバル座標
                global_sx = x + sx
                global_sy = y + sy
                
                # v40 fix: 塗りつぶしと同じmarginで拡張して切り出し（枠線残り防止）
                margin = 7  # v40: 5から7に拡張
                crop_sx = max(0, global_sx - margin)
                crop_sy = max(0, global_sy - margin)
                crop_ex = min(img.shape[1], global_sx + sw + margin)
                crop_ey = min(img.shape[0], global_sy + sh + margin)
                
                # オリジナル画像からサブ要素を切り出し（塗りつぶし前の状態、margin分拡張）
                # 注意：elementは既に塗りつぶし済みなので、imgから直接取得
                sub_crop = img[crop_sy:crop_ey, crop_sx:crop_ex].copy()
                
                # 拡張後のサイズを計算
                expanded_w = crop_ex - crop_sx
                expanded_h = crop_ey - crop_sy
                
                # v41: テキスト重複チェックを復活（v40では無効化されていた）
                is_text_overlap = is_text_region((global_sx, global_sy, sw, sh), text_bboxes_pixel, threshold=0.2)
                
                if is_text_overlap:
                    # v41: 重複しているテキストボックスの色をtext_bboxes_with_colorsから取得
                    def get_overlapping_text_colors_from_data(sub_bbox, text_bboxes_with_colors):
                        """重複しているテキストボックスの色をデータから取得（画像サンプリングではなく）"""
                        if not text_bboxes_with_colors:
                            return []
                        
                        sx, sy, sw, sh = sub_bbox
                        sub_right, sub_bottom = sx + sw, sy + sh
                        text_colors = []
                        
                        for bbox, color_rgb in text_bboxes_with_colors:
                            tx1, ty1, tx2, ty2 = bbox
                            
                            # 重なり判定
                            overlap_x1 = max(sx, tx1)
                            overlap_y1 = max(sy, ty1)
                            overlap_x2 = min(sub_right, tx2)
                            overlap_y2 = min(sub_bottom, ty2)
                            
                            if overlap_x2 > overlap_x1 and overlap_y2 > overlap_y1:
                                # color_rgbはRGB形式で渡されるため、BGRに変換
                                color_bgr = (color_rgb[2], color_rgb[1], color_rgb[0]) if color_rgb else (0, 0, 0)
                                text_colors.append(color_bgr)
                        
                        return text_colors
                    
                    overlapping_text_colors = get_overlapping_text_colors_from_data(
                        (global_sx, global_sy, sw, sh), text_bboxes_with_colors
                    )
                    
                    # テキスト色と親の塗りつぶし色の差を計算
                    exempt_due_to_color = False
                    for text_color in overlapping_text_colors:
                        # BGR差を計算
                        color_diff = abs(int(text_color[0]) - int(parent_fill_color[0])) + \
                                     abs(int(text_color[1]) - int(parent_fill_color[1])) + \
                                     abs(int(text_color[2]) - int(parent_fill_color[2]))
                        
                        if color_diff <= 70:  # v41: 閾値70
                            exempt_due_to_color = True
                            print(f"    [v41] Text color BGR={text_color}, parent_fill BGR={parent_fill_color}, diff={color_diff} <= 70 -> EXEMPT")
                            break
                    
                    if exempt_due_to_color:
                        # テキスト色が親と近い → 免除（保持）
                        print(f"    [v41] Sub at ({crop_sx},{crop_sy},{expanded_w},{expanded_h}): text overlap=True, but COLOR EXEMPT")
                    else:
                        # テキスト色が親と遠い → 削除（continue）
                        if overlapping_text_colors:
                             first_color = overlapping_text_colors[0]
                             diff = abs(int(first_color[0]) - int(parent_fill_color[0])) + \
                                    abs(int(first_color[1]) - int(parent_fill_color[1])) + \
                                    abs(int(first_color[2]) - int(parent_fill_color[2]))
                             print(f"    [v41-DEBUG] Skipping: Text BGR={first_color}, Parent BGR={parent_fill_color}, Diff={diff} > 70")
                        
                        print(f"    [v41] Sub at ({crop_sx},{crop_sy},{expanded_w},{expanded_h}): text overlap=True, no color exempt, skipping")
                        continue
                else:
                    print(f"    [v41] Sub at ({crop_sx},{crop_sy},{expanded_w},{expanded_h}): text overlap=False, adding to overlay (expanded by margin={margin})")
                
                # サブ要素を保存
                sub_elem_path = os.path.join(output_dir, f"sub_element_{i}_{sx}_{sy}.png")
                
                # RGBA変換（不透明）
                sub_rgba = cv2.cvtColor(sub_crop, cv2.COLOR_BGR2BGRA)
                sub_rgba[:, :, 3] = 255  # 完全不透明
                cv2.imwrite(sub_elem_path, sub_rgba)
                
                # elements listに追加（親の後に描画されるよう、後で追加）
                # v40 fix: bboxも拡張後の座標・サイズに更新
                elements.append({
                    "path": sub_elem_path,
                    "bbox": (crop_sx, crop_sy, expanded_w, expanded_h),
                    "area": expanded_w * expanded_h,
                    "is_sub_overlay": True  # v37: サブオーバーレイフラグ
                })
                print(f"    [v37] Added sub-element overlay: {sub_elem_path}")
            
            # v17: NESTED処理した親要素全体を不透明にする
            # 背景色に近い親要素でも透過しないようにする
            element_mask_crop[:, :] = 255
            print(f"    [NESTED] v17: Made entire parent element opaque")
            
            nested_detect_count += 1


        # テキスト領域と20%以上重なる場合はスキップ
        is_overlapping = is_text_region((x, y, w, h), text_bboxes_pixel, threshold=0.2)
        
        if is_overlapping:
            text_overlap_count += 1
            continue

        # RGBA変換（透過対応）
        element_rgba = cv2.cvtColor(element, cv2.COLOR_BGR2BGRA)
        element_rgba[:, :, 3] = element_mask_crop
        
        elem_path = os.path.join(output_dir, f"element_{i}.png")
        cv2.imwrite(elem_path, element_rgba)
        
        elements.append({
            "path": elem_path,
            "bbox": (x, y, w, h),
            "area": area
        })
    
    print(f"Excluded {text_overlap_count} text-overlapping elements")
    print(f"Nested detection on {nested_detect_count} graphic elements")
    print(f"Recovered {len(recovered_texts)} text elements from graphics")
    print(f"Kept {len(elements)} graphic elements")
    
    return elements, recovered_texts

def add_slide_from_page(prs, page, page_image_path, page_width, page_height):
    """既存のプレゼンテーションにスライドを追加"""
    print("  [DEBUG] Adding blank slide...")
    # 空白スライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 背景色を設定
    page_img = Image.open(page_image_path)
    bg_color = get_background_color(page_img)
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
    
    # JSON補正データの確認 (replace_all対応)
    corrections = load_ocr_corrections()
    page_key = f"page_{page.number + 1}"
    force_ocr = False
    
    if page_key in corrections and corrections[page_key].get("replace_all"):
        force_ocr = True
        print(f"  [DEBUG] Force replacing text blocks for {page_key} because replace_all=True")

    # テキストブロックを抽出
    text_blocks = []
    if not force_ocr:
        text_blocks = extract_text_blocks_native(page)
    
    text_bboxes_pixel = []
    
    if len(text_blocks) == 0 or force_ocr:
        print("  [DEBUG] Using OCR/JSON text blocks...")
        text_blocks, text_bboxes_pixel = extract_text_with_ocr_word_level(page_image_path, page_width, page_height)
        print("  [DEBUG] OCR finished.")
    
    # 画像読み込み（色検出用）
    cv_img = read_image_cv2(page_image_path)
    img_width, img_height = page_img.size
    
    # まずグラフィック要素を追加（背景側）
    print("  [DEBUG] Detecting graphic elements...")
    # OCRワードフィルタリングを無効化（矢印が消える問題を回避）
    # all_text_bboxes = list(text_bboxes_pixel)
    # for word in OCR_ALL_WORDS:
    #     word_bbox = [word['left'], word['top'], word['right'], word['bottom']]
    #     all_text_bboxes.append(word_bbox)
    
    # v39: text_blocksから色情報付きbboxリストを作成
    text_bboxes_with_colors = []
    for block in text_blocks:
        # pixel座標を取得（既にtext_bboxes_pixelに含まれているはずだが、色と対応付けるために再計算）
        bbox_pt = block.get("bbox", [0, 0, 0, 0])
        bbox_pixel = [
            int(bbox_pt[0] / page_width * img_width),
            int(bbox_pt[1] / page_height * img_height),
            int(bbox_pt[2] / page_width * img_width),
            int(bbox_pt[3] / page_height * img_height)
        ]
        color = block.get("color", [0, 0, 0])  # RGB形式
        text_bboxes_with_colors.append((bbox_pixel, color))
    
    # NOTE: Threshold logic is inside detect_graphic_elements. Ideally we pass it as arg, but hardcoded for now.
    # We rely on the fix applied to detect_graphic_elements above (threshold=0.05).
    elements, recovered_texts_from_graphics = detect_graphic_elements(page_image_path, OUTPUT_DIR, text_bboxes_pixel, text_bboxes_with_colors)
    print(f"  [DEBUG] Found {len(elements)} graphic elements")
    
    # リカバリーされたテキストを座標変換してtext_blocksに追加
    for r_block in recovered_texts_from_graphics:
        # pixel bbox -> pt bbox
        px_bbox = r_block["bbox"]
        pt_bbox = [
            px_bbox[0] / img_width * page_width,
            px_bbox[1] / img_height * page_height,
            px_bbox[2] / img_width * page_width,
            px_bbox[3] / img_height * page_height
        ]
        
        # NFKC正規化と丸数字置換（ここでも適用）
        text = unicodedata.normalize('NFKC', r_block["text"])
        char_map = {'①':'1', '②':'2', '③':'3', '④':'4'} # リカバリー分は念のためマップも適用
        for k, v in char_map.items():
            text = text.replace(k, v)
        
        # フォントサイズ推定（高さの80%）
        font_size = (pt_bbox[3] - pt_bbox[1]) * 0.8
        
        text_blocks.append({
            "text": text,
            "bbox": pt_bbox,
            "font_size": font_size,
            "color": [0,0,0], # 黒と仮定
            "word_details": [] # 詳細なし
        })

    # v40 fix: サブオーバーレイと通常要素を分離して処理
    # 通常要素は面積順で上位50個、サブオーバーレイは全て追加
    regular_elements = [e for e in elements if not e.get("is_sub_overlay", False)]
    sub_overlay_elements = [e for e in elements if e.get("is_sub_overlay", False)]
    
    print(f"  [v40] Regular elements: {len(regular_elements)}, Sub-overlay elements: {len(sub_overlay_elements)}")
    
    # 通常要素を追加（背景側、大きい順に50個まで）
    for elem in sorted(regular_elements, key=lambda x: x["area"], reverse=True)[:50]:
        x, y, w, h = elem["bbox"]
        
        elem_left_pt = x / img_width * page_width
        elem_top_pt = y / img_height * page_height
        elem_width_pt = w / img_width * page_width
        elem_height_pt = h / img_height * page_height
        
        try:
            slide.shapes.add_picture(elem["path"], Pt(elem_left_pt), Pt(elem_top_pt), Pt(elem_width_pt), Pt(elem_height_pt))
        except Exception as e:
            print(f"  [ERROR] Failed to add regular picture: {e}")
    
    # サブオーバーレイ要素を追加（前面側、全て追加）
    for elem in sub_overlay_elements:
        x, y, w, h = elem["bbox"]
        
        elem_left_pt = x / img_width * page_width
        elem_top_pt = y / img_height * page_height
        elem_width_pt = w / img_width * page_width
        elem_height_pt = h / img_height * page_height
        
        try:
            slide.shapes.add_picture(elem["path"], Pt(elem_left_pt), Pt(elem_top_pt), Pt(elem_width_pt), Pt(elem_height_pt))
            print(f"  [v40] Added sub-overlay at ({x},{y},{w},{h})")
        except Exception as e:
            print(f"  [ERROR] Failed to add sub-overlay picture: {e}")
    
    # テキストを追加（前面側）
    print(f"  [DEBUG] Adding {len(text_blocks)} text blocks...")
    
    # ★ v43: SSIMフォントサイズを事前計算し、JSONのfont_size_ptが同じグループ内で最大サイズに統一
    # Step 1: 各ブロックのSSIMサイズを事前計算
    ssim_sizes_cache = {}  # block index -> ssim_font_size
    json_size_to_ssim_sizes = {}  # json_font_size_pt -> [ssim_size1, ssim_size2, ...]
    
    print("  [v43] Pre-computing SSIM font sizes for unification...")
    for idx, block in enumerate(text_blocks):
        # word_detailsがないブロック（JSONブロック）のみ対象
        if "word_details" not in block or not block["word_details"]:
            bbox = block["bbox"]
            json_img_bbox = [
                int(bbox[0] / page_width * img_width),
                int(bbox[1] / page_height * img_height),
                int(bbox[2] / page_width * img_width),
                int(bbox[3] / page_height * img_height)
            ]
            
            json_font_family = block.get("font_family", "Noto Sans JP")
            json_is_bold = block.get("is_bold", False)
            json_font_size = block.get("font_size_pt", block.get("font_size", 12.0))
            
            if cv_img is not None:
                # OCR bboxを検索
                ocr_bbox = find_ocr_bbox_for_text(block["text"], OCR_ALL_WORDS, json_img_bbox)
                img_bbox = ocr_bbox if ocr_bbox else json_img_bbox
                
                font_props = detect_font_properties_v3(
                    block["text"], 
                    img_bbox, 
                    cv_img,
                    json_font_family=json_font_family,
                    json_font_size_pt=json_font_size,
                    json_is_bold=json_is_bold,
                    debug=False  # 事前計算なのでデバッグ抑制
                )
                ssim_size = font_props["font_size"]
            else:
                ssim_size = json_font_size
            
            ssim_sizes_cache[idx] = ssim_size
            
            # JSONのfont_size_ptでグループ化
            json_size_key = round(json_font_size, 1)  # 小数点1桁で丸めてグループ化
            if json_size_key not in json_size_to_ssim_sizes:
                json_size_to_ssim_sizes[json_size_key] = []
            json_size_to_ssim_sizes[json_size_key].append((idx, ssim_size))
    
    # Step 2: 各グループ内で最大サイズを決定
    unified_sizes = {}  # block index -> unified_font_size
    for json_size_key, ssim_list in json_size_to_ssim_sizes.items():
        max_ssim_size = max(s for _, s in ssim_list)
        print(f"    [v43] JSON size {json_size_key}pt: {len(ssim_list)} blocks -> unified to {max_ssim_size:.1f}pt")
        for idx, _ in ssim_list:
            unified_sizes[idx] = max_ssim_size
    
    print(f"  [v43] Unified {len(unified_sizes)} blocks across {len(json_size_to_ssim_sizes)} size groups")

    for i, block in enumerate(text_blocks):
        if i % 10 == 0: print(f"    [DEBUG] Text block {i}/{len(text_blocks)}")
        bbox = block["bbox"]
        
        if "word_details" in block and block["word_details"]:
            valid_bboxes = [w['bbox_pixel'] for w in block["word_details"] if w['bbox_pixel']]
            if valid_bboxes:
                all_x1 = min([b[0] for b in valid_bboxes])
                all_y1 = min([b[1] for b in valid_bboxes])
                all_x2 = max([b[2] for b in valid_bboxes])
                all_y2 = max([b[3] for b in valid_bboxes])
                
                pixel_bbox = [all_x1, all_y1, all_x2, all_y2]
                refined_pixel_bbox = refine_bbox_vertical(cv_img, pixel_bbox, bg_color)
                
                new_y1_pt = refined_pixel_bbox[1] / img_height * page_height
                new_y2_pt = refined_pixel_bbox[3] / img_height * page_height
                
                bbox[1] = new_y1_pt
                bbox[3] = new_y2_pt
        
        left = Pt(bbox[0])
        top = Pt(bbox[1])
        width = Pt((bbox[2] - bbox[0]))  # 幅調整なし（正確なBBoxを信頼）
        height = Pt(bbox[3] - bbox[1])
        
        if width < Pt(10): width = Pt(20) # 最小値
        if height < Pt(10): height = Pt(20)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False # 折り返しなし
        tf.auto_size = None # 固定サイズ
        
        tf.margin_left = 0
        tf.margin_top = 0
        tf.margin_right = 0
        tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        
        # v5: フォントサイズの高度な判定
        if "font_size_fixed" in block:
            final_font_size_pt = block["font_size_fixed"]
        else:
            # 画像のピクセル高さから直接推定 (より正確)
            if "word_details" in block and block["word_details"]:
                valid_bboxes = [w['bbox_pixel'] for w in block["word_details"] if w['bbox_pixel']]
                if valid_bboxes:
                    px_h = max([b[3] for b in valid_bboxes]) - min([b[1] for b in valid_bboxes])
                    final_font_size_pt = estimate_font_size_v5(px_h, img_height, page_height)
                else:
                    final_font_size_pt = block["font_size"]
            else:
                # word_detailsがない場合はBBox(pt)から逆算
                final_font_size_pt = (bbox[3] - bbox[1]) * 0.92
        # 背景色設定 - ユーザー要望により一時無効化
        # if "background_color" in block:
        #    bg_rgb = block["background_color"]
        #    textbox.fill.solid()
        #    textbox.fill.fore_color.rgb = RGBColor(int(bg_rgb[0]), int(bg_rgb[1]), int(bg_rgb[2]))
        
        # フォント設定（明示的）を使用
        # ...
        if "word_details" in block and block["word_details"] and cv_img is not None:
            p.text = ""
            p.line_spacing = 1.0
            
            analyzed_words = []
            current_color_rgb = None
            
            for word_info in block.get("word_details", []):
                word_text = word_info['text']
                word_bbox = word_info['bbox_pixel']
                
                color_rgb = (0, 0, 0)
                weight_ratio = 0.0
                stroke_var = 0.0
                
                if word_bbox is None:
                    if current_color_rgb is not None:
                        color_rgb = current_color_rgb
                else:
                    color_rgb = detect_text_color(cv_img, word_bbox, bg_color)
                    current_color_rgb = color_rgb
                    weight_ratio, stroke_var = detect_font_weight_v5(cv_img, word_bbox, bg_color)
                
                analyzed_words.append({
                    'text': word_text,
                    'color': color_rgb,
                    'ratio': weight_ratio,
                    'var': stroke_var
                })
            
            # 色の統一処理
            color_counts = {}
            for aw in analyzed_words:
                c = aw['color']
                color_counts[c] = color_counts.get(c, 0) + len(aw['text'])
            
            sorted_colors = sorted(color_counts.keys(), key=lambda k: color_counts[k], reverse=True)
            
            color_map = {}
            merged_colors = []
            
            for c in sorted_colors:
                found_match = False
                for mc in merged_colors:
                    dist = ((c[0]-mc[0])**2 + (c[1]-mc[1])**2 + (c[2]-mc[2])**2)**0.5
                    if dist < 40:
                        color_map[c] = mc
                        found_match = True
                        break
                
                if not found_match:
                    merged_colors.append(c)
                    color_map[c] = c
            
            for aw in analyzed_words:
                aw['color'] = color_map[aw['color']]
            
            runs_groups = []
            if analyzed_words:
                current_run_group = {'text': "", 'color': analyzed_words[0]['color'], 'ratios': [], 'vars': []}
                
                for aw in analyzed_words:
                    if aw['color'] != current_run_group['color']:
                        runs_groups.append(current_run_group)
                        current_run_group = {'text': "", 'color': aw['color'], 'ratios': [], 'vars': []}
                    
                    current_run_group['text'] += aw['text']
                    if aw['ratio'] > 0:
                        current_run_group['ratios'].append(aw['ratio'])
                        current_run_group['vars'].append(aw['var'])
                
                runs_groups.append(current_run_group)
            
            for i, group in enumerate(runs_groups):
                if i == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                else:
                    run = p.add_run()
                
                run.text = group['text']
                run.font.size = Pt(final_font_size_pt)
                
                c = group['color']
                run.font.color.rgb = RGBColor(c[0], c[1], c[2])
                
                ratios = group['ratios']
                vars_list = group['vars']
                
                median_ratio = sorted(ratios)[len(ratios) // 2] if ratios else 0
                median_var = sorted(vars_list)[len(vars_list) // 2] if vars_list else 0
                
                selected_font, is_bold_detected = get_font_for_text_v5(median_ratio, median_var, c)
                
                # JSON指定があれば優先
                if "font_family" in block:
                    set_font_for_run(run, block["font_family"])
                else:
                    run.font.name = selected_font
                
                # JSON指定があれば優先、なければ検出結果
                if "is_bold" in block:
                    if block["is_bold"]:
                        run.font.bold = True
                elif is_bold_detected:
                    run.font.bold = True
        else:
            # ★ JSONブロック（word_detailsがない）の処理 - SSIMベースのフォント検出
            # v28: colors配列対応 - 複数色の場合はrunを分割
            
            # JSON指定のbboxを画像座標に変換
            json_img_bbox = [
                int(bbox[0] / page_width * img_width),
                int(bbox[1] / page_height * img_height),
                int(bbox[2] / page_width * img_width),
                int(bbox[3] / page_height * img_height)
            ]
            
            # OCRワードから正確なbboxを検索
            ocr_bbox = find_ocr_bbox_for_text(block["text"], OCR_ALL_WORDS, json_img_bbox)
            
            if ocr_bbox:
                img_bbox = ocr_bbox
                print(f"  [JSON Block] '{block['text'][:15]}...'")
                print(f"    OCR bbox found: {img_bbox}")
                use_width_based_size = False
            else:
                img_bbox = json_img_bbox
                print(f"  [JSON Block] '{block['text'][:15]}...'")
                print(f"    Using JSON bbox (no OCR match): {img_bbox}")
                use_width_based_size = True
            
            # SSIMベースのサイズ検出
            json_font_family = block.get("font_family", "Noto Sans JP")
            json_is_bold = block.get("is_bold", False)
            json_font_size = block.get("font_size_pt", block.get("font_size", 12.0))
            
            if cv_img is not None:
                font_props = detect_font_properties_v3(
                    block["text"], 
                    img_bbox, 
                    cv_img,
                    json_font_family=json_font_family,
                    json_font_size_pt=json_font_size,
                    json_is_bold=json_is_bold,
                    debug=True
                )
                
                # フォントサイズ決定 - v43: 統一サイズを優先使用
                if "font_size_fixed" in block:
                    final_font_size_pt = block["font_size_fixed"]
                    print(f"    [SSIM] Overridden by JSON: {final_font_size_pt}pt")
                elif i in unified_sizes:
                    # v43: 統一サイズを使用
                    final_font_size_pt = unified_sizes[i]
                    original_ssim = ssim_sizes_cache.get(i, 0)
                    print(f"    [v43] Unified size: {original_ssim:.1f}pt -> {final_font_size_pt:.1f}pt")
                elif "font_size_pt" in block:
                    # JSONで指定されたfont_size_ptを使用
                    final_font_size_pt = block["font_size_pt"]
                    print(f"    [JSON] Using font_size_pt: {final_font_size_pt}pt")
                else:
                    # SSIMの値を直接使用（WIDTH blending無効化）
                    final_font_size_pt = font_props["font_size"]
                    print(f"    [SSIM] Detected size: {final_font_size_pt:.1f}pt")
            else:
                final_font_size_pt = json_font_size
            
            # v28: colors配列対応 - 複数のrunに分割
            text_content = block["text"]
            colors_array = block.get("multi_colors", [])  # JSONのcolorsはmulti_colorsとして保存される
            
            # Debug: colors配列の確認
            print(f"    [v28 DEBUG] multi_colors = {colors_array}")
            print(f"    [v28 DEBUG] block keys = {list(block.keys())}")
            
            # v31: SSIMの太字判定結果を取得（cv_imgがある場合）
            ssim_is_bold = font_props.get("is_bold", json_is_bold) if cv_img is not None else json_is_bold
            final_is_bold = ssim_is_bold  # SSIMの結果を優先
            print(f"    [v31] Bold: JSON={json_is_bold}, SSIM={ssim_is_bold} -> Final={final_is_bold}")
            
            if colors_array and len(colors_array) > 0:
                # colors配列がある場合：rangeに従ってrunを分割
                p.text = ""  # 一旦クリア
                
                # v29: colors rangeがテキスト全体をカバーしているか確認し、不足分を補完
                covered_end = 0  # 現在までにカバーされた最後のインデックス
                
                for color_idx, color_info in enumerate(colors_array):
                    color_range = color_info.get("range", [0, len(text_content)])
                    color_rgb = color_info.get("rgb", [0, 0, 0])
                    
                    start_idx = color_range[0]
                    end_idx = color_range[1]
                    
                    # v29: rangeが不連続の場合、間のテキストをデフォルト色で追加
                    if start_idx > covered_end:
                        gap_text = text_content[covered_end:start_idx]
                        if gap_text:
                            run = p.add_run()
                            run.text = gap_text
                            run.font.size = Pt(final_font_size_pt)
                            run.font.color.rgb = RGBColor(0, 0, 0)  # デフォルト黒
                            set_font_for_run(run, json_font_family)
                            run.font.bold = final_is_bold
                            print(f"    [v29 GAP] '{gap_text}' (default black)")
                    
                    # テキストの該当部分を抽出
                    run_text = text_content[start_idx:end_idx]
                    
                    if run_text:
                        run = p.add_run()
                        run.text = run_text
                        run.font.size = Pt(final_font_size_pt)
                        run.font.color.rgb = RGBColor(int(color_rgb[0]), int(color_rgb[1]), int(color_rgb[2]))
                        set_font_for_run(run, json_font_family)
                        run.font.bold = final_is_bold
                        print(f"    [v29 MULTI-COLOR] Run: '{run_text}' color=RGB{tuple(color_rgb)}")
                    
                    covered_end = max(covered_end, end_idx)
                
                # v29: 残りのテキストがあれば追加（最後のrangeがテキスト末尾まで届いていない場合）
                if covered_end < len(text_content):
                    remaining_text = text_content[covered_end:]
                    if remaining_text:
                        run = p.add_run()
                        run.text = remaining_text
                        run.font.size = Pt(final_font_size_pt)
                        # 最後の色を引き継ぐか、デフォルト黒
                        last_color = colors_array[-1].get("rgb", [0, 0, 0]) if colors_array else [0, 0, 0]
                        run.font.color.rgb = RGBColor(int(last_color[0]), int(last_color[1]), int(last_color[2]))
                        set_font_for_run(run, json_font_family)
                        run.font.bold = final_is_bold
                        print(f"    [v29 REMAINING] '{remaining_text}' color=RGB{tuple(last_color)}")
                
                print(f"    [v29] Multi-color processing complete")
            else:
                # colors配列がない場合：従来通り単一run
                p.text = text_content
                if p.runs:
                    run = p.runs[0]
                    run.font.size = Pt(final_font_size_pt)
                    
                    # 色設定（旧color形式または画像検出）
                    if "color" in block:
                        color = block["color"]
                        color_rgb = (int(color[0]), int(color[1]), int(color[2]))
                    elif cv_img is not None:
                        color_rgb = detect_text_color(cv_img, img_bbox, bg_color)
                    else:
                        color_rgb = (0, 0, 0)
                    
                    run.font.color.rgb = RGBColor(color_rgb[0], color_rgb[1], color_rgb[2])
                    print(f"    [SSIM] Color: {color_rgb}")
                    
                    # フォント設定
                    set_font_for_run(run, json_font_family)
                    print(f"    [SSIM] Font applied: {json_font_family}")
                    
                    # 太字設定
                    run.font.bold = final_is_bold
                    print(f"    [SSIM] Bold: {run.font.bold}")



def create_pptx(page, page_image_path, output_path, page_width, page_height):
    """PowerPointを生成"""
    # プレゼンテーション作成
    prs = Presentation()
    
    # スライドサイズをPDFページサイズに合わせる
    # 1 pt = 12700 Emu
    slide_width_emu = int(page_width * 12700)
    slide_height_emu = int(page_height * 12700)
    
    prs.slide_width = Emu(slide_width_emu)
    prs.slide_height = Emu(slide_height_emu)
    print(f"Slide size set to: {page_width}x{page_height} pt ({slide_width_emu}x{slide_height_emu} Emu)")
    
    # 利用可能なフォントのデバッグ（Windows）
    try:
        from matplotlib import font_manager
        fonts = font_manager.findSystemFonts(fontpaths=None, fontext='ttf')
        print(f"System fonts found: {len(fonts)}")
        # 游ゴシック、メイリオ等があるかチェック
        jp_fonts = [f for f in fonts if "YuGo" in f or "Meiryo" in f or "HGP" in f]
        print(f"Japanese fonts candidates: {[os.path.basename(f) for f in jp_fonts]}")
    except:
        print("Could not list system fonts")

    # 空白スライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 背景色を設定
    page_img = Image.open(page_image_path)
    bg_color = get_background_color(page_img)
    print(f"Background color (RGB): {bg_color}")
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
    
    # JSON補正データの確認 (replace_all対応)
    corrections = load_ocr_corrections()
    page_key = f"page_{page.number + 1}"
    corrections = load_ocr_corrections()
    page_key = f"page_{page.number + 1}"
    print(f"Current page key: {page_key}")
    if page_key in corrections:
        print(f"  Found in corrections. replace_all={corrections[page_key].get('replace_all')}")
    else:
        print(f"  Not found in corrections.")
    force_ocr = False
    
    if page_key in corrections and corrections[page_key].get("replace_all"):
        force_ocr = True
        print(f"Force replacing text blocks for {page_key} because replace_all=True")
    
    # テキストブロックを抽出
    text_blocks = []
    if not force_ocr:
        text_blocks = extract_text_blocks_native(page)
        print(f"Native text blocks: {len(text_blocks)}")
    
    if len(text_blocks) == 0 or force_ocr:
        print("Using OCR/JSON text blocks...")
        text_blocks, text_bboxes_pixel = extract_text_with_ocr_word_level(page_image_path, page_width, page_height)
        print(f"OCR/JSON text blocks: {len(text_blocks)}")
    
    # 画像読み込み（色検出用）
    cv_img = read_image_cv2(page_image_path)
    img_width, img_height = page_img.size
    
    # まずグラフィック要素を追加（背景側）
    elements = detect_graphic_elements(page_image_path, OUTPUT_DIR, text_bboxes_pixel)
    print(f"Adding {len(elements)} graphic elements")
    
    for elem in sorted(elements, key=lambda x: x["area"], reverse=True)[:50]:
        x, y, w, h = elem["bbox"]
        
        # ピクセル座標 -> PDFポイント座標 -> Emu
        # x / img_width * page_width -> pt
        elem_left_pt = x / img_width * page_width
        elem_top_pt = y / img_height * page_height
        elem_width_pt = w / img_width * page_width
        elem_height_pt = h / img_height * page_height
        
        try:
            slide.shapes.add_picture(elem["path"], Pt(elem_left_pt), Pt(elem_top_pt), Pt(elem_width_pt), Pt(elem_height_pt))
        except Exception as e:
            print(f"Failed to add element: {e}")
    
    # テキストを追加（前面側）
    for block in text_blocks:
        bbox = block["bbox"]
        
        # 垂直位置の精密補正（bboxはPDFポイント座標だが、画像はピクセルなので変換が必要）
        # pixel_bbox = [x * scale_x, y * scale_y, ...]
        # ここでは簡易的に、現在のblock["word_details"]からピクセルbboxを再構成するか、
        # あるいはblock自体にpixel_bbox情報を保持させておくのが正道だが、
        # extract_text_blocks_nativeはnative座標しか返さない。
        # word_detailsがあればそれを使う。
        
        if "word_details" in block and block["word_details"]:
             # ワード全体のピクセル外接矩形を計算
             all_x1 = min([w['bbox_pixel'][0] for w in block["word_details"] if w['bbox_pixel']])
             all_y1 = min([w['bbox_pixel'][1] for w in block["word_details"] if w['bbox_pixel']])
             all_x2 = max([w['bbox_pixel'][2] for w in block["word_details"] if w['bbox_pixel']])
             all_y2 = max([w['bbox_pixel'][3] for w in block["word_details"] if w['bbox_pixel']])
             
             pixel_bbox = [all_x1, all_y1, all_x2, all_y2]
             refined_pixel_bbox = refine_bbox_vertical(cv_img, pixel_bbox, bg_color)
             
             # ピクセル座標からポイント座標へ逆変換
             # pt = pixel / img_width * page_width
             new_y1_pt = refined_pixel_bbox[1] / img_height * page_height
             new_y2_pt = refined_pixel_bbox[3] / img_height * page_height
             
             # 高さを更新
             bbox[1] = new_y1_pt
             bbox[3] = new_y2_pt
        
        # bboxはPDFポイント座標
        left = Pt(bbox[0])
        top = Pt(bbox[1])
        width = Pt(bbox[2] - bbox[0])
        height = Pt(bbox[3] - bbox[1])
        
        # 最小サイズの保護
        if width < Pt(20): width = Pt(50)
        if height < Pt(10): height = Pt(20)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False
        # マージンをゼロにして配置ズレを解消
        tf.margin_left = 0
        tf.margin_top = 0
        tf.margin_right = 0
        tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        
        # 基本フォントサイズ（高さ基準）
        base_font_size_pt = block["font_size"]
        
        # フォントサイズ自動調整（JSON指定があればそれを優先）
        text_content = block["text"].replace(" ", "")  # スペースを除去
        box_width_pt = bbox[2] - bbox[0]
        final_font_size_pt = base_font_size_pt
        
        # JSONで指定されたフォントサイズがあればそれを使う
        if "font_size_fixed" in block:
            final_font_size_pt = block["font_size_fixed"]
            # 自動計算をスキップするためのフラグとして利用
            text_content = "" # 計算スキップ
            print(f"  [DEBUG] Fixed font size for '{block['text'][:10]}...': {final_font_size_pt}")
        else:
            print(f"  [DEBUG] No fixed font size for '{block['text'][:10]}...'. Base: {base_font_size_pt}")
        
        if len(text_content) > 0 and box_width_pt > 0:
            try:
                # HGPゴシックE（PPTで使用するフォント）で計測
                font_path = r"C:\Windows\Fonts\HGRGE.TTC"
                base_size = 100
                font = ImageFont.truetype(font_path, base_size)
                measured_width = font.getlength(text_content)
                
                if measured_width > 0:
                    # PPTとPillowの描画差異を吸収する補正係数
                    # Pillowの計測幅がPPTより大きいため、係数を上げて補正
                    ppt_pillow_ratio = 1.13
                    
                    calculated_size = base_size * (box_width_pt / measured_width) * ppt_pillow_ratio
                    
                    # 極端な値を制限
                    final_font_size_pt = max(5, min(calculated_size, base_font_size_pt * 2))
                    
                    print(f"HGP Fit: '{text_content[:10]}...' Box={box_width_pt:.0f} Meas={measured_width:.0f} -> Size={final_font_size_pt:.1f}pt")
            except Exception as e:
                print(f"Font measurement error: {e}")
                # フォールバック: 元のサイズを使用
                final_font_size_pt = base_font_size_pt

        # 単語ごとにrunを作成して色を設定
        if "word_details" in block and cv_img is not None:
            # 段落のテキストをクリア
            p.text = ""
            # 行間を1.0行（シングル）に設定して垂直位置を安定させる
            p.line_spacing = 1.0 # 1.0 = Single spacing
            
            # 1. まず全単語の色と太さ情報を計算してリスト化
            analyzed_words = []
            
            current_color_rgb = None
            if "word_details" in block:
                for word_info in block["word_details"]:
                    word_text = word_info['text']
                    word_bbox = word_info['bbox_pixel']
                    
                    color_rgb = (0, 0, 0) # Default
                    weight_ratio = 0.0
                    stroke_var = 0.0
                    
                    if word_bbox is None:
                        # スペースの場合：前の色を引き継ぐ（未設定なら黒）
                        if current_color_rgb is not None:
                            color_rgb = current_color_rgb
                    else:
                        color_rgb = detect_text_color(cv_img, word_bbox, bg_color)
                        current_color_rgb = color_rgb
                        # 太さ比率と変動係数を取得
                        weight_ratio, stroke_var = detect_font_weight_v5(cv_img, word_bbox, bg_color)
                
                    analyzed_words.append({
                        'text': word_text,
                        'color': color_rgb,
                        'ratio': weight_ratio,
                        'var': stroke_var
                    })
            
            # 色の統一処理（クラスタリング）
            # 1. 出現色のカウント
            color_counts = {}
            for aw in analyzed_words:
                c = aw['color']
                color_counts[c] = color_counts.get(c, 0) + len(aw['text'])
            
            # 2. 頻度順にソート（最も使われている色が代表色の候補）
            sorted_colors = sorted(color_counts.keys(), key=lambda k: color_counts[k], reverse=True)
            
            # 3. 近似色の統合マップ作成
            color_map = {}
            merged_colors = []
            
            for c in sorted_colors:
                # 既に統合済みの代表色の中に近いものがあるか探す
                found_match = False
                for mc in merged_colors:
                    # RGB距離計算
                    dist = ((c[0]-mc[0])**2 + (c[1]-mc[1])**2 + (c[2]-mc[2])**2)**0.5
                    if dist < 40: # 閾値（要調整: 30-50くらい）
                        color_map[c] = mc
                        found_match = True
                        break
                
                if not found_match:
                    merged_colors.append(c)
                    color_map[c] = c
            
            # 4. 色の置換
            for aw in analyzed_words:
                aw['color'] = color_map[aw['color']]

            # 2. 色ごとにRunをグループ化
            runs_groups = []
            if analyzed_words:
                current_run_group = {'text': "", 'color': analyzed_words[0]['color'], 'ratios': [], 'vars': []}
                
                for aw in analyzed_words:
                    if aw['color'] != current_run_group['color']:
                        # 色が変わったら今のグループを保存して新しいグループ開始
                        runs_groups.append(current_run_group)
                        current_run_group = {'text': "", 'color': aw['color'], 'ratios': [], 'vars': []}
                    
                    current_run_group['text'] += aw['text']
                    if aw['ratio'] > 0: # スペース等は除外
                         current_run_group['ratios'].append(aw['ratio'])
                         current_run_group['vars'].append(aw['var'])
                
                runs_groups.append(current_run_group)
            
            # 3. 各Runを追加
            for i, group in enumerate(runs_groups):
                if i == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                else:
                    run = p.add_run()
                
                run.text = group['text']
                run.font.size = Pt(final_font_size_pt)
                
                # 色設定
                c = group['color']
                run.font.color.rgb = RGBColor(c[0], c[1], c[2])
                
                # 太字判定＆明朝判定：グループ内の中央値を使用
                ratios = group['ratios']
                vars_list = group['vars']
                
                median_ratio = 0
                median_var = 0
                if ratios:
                    median_ratio = sorted(ratios)[len(ratios) // 2]
                if vars_list:
                    median_var = sorted(vars_list)[len(vars_list) // 2]
                
                selected_font, is_bold = get_font_for_text_v5(median_ratio, median_var, c)
                
                print(f"Run '{group['text'][:5]}...' Ratio:{median_ratio:.3f} Var:{median_var:.3f} -> {selected_font} (Bold:{is_bold})")
                
                run.font.name = selected_font
                if is_bold:
                    run.font.bold = True
            
            print(f"Text block processed: {len(runs_groups)} color runs created")
        # JSONブロック（word_detailsがない場合）のスタイル適用ロジック
        else:
            p.text = block["text"]
            if p.runs:
                run = p.runs[0]
                
                # 画像座標に変換
                img_bbox = [
                    int(bbox[0] / page_width * img_width),
                    int(bbox[1] / page_height * img_height),
                    int(bbox[2] / page_width * img_width),
                    int(bbox[3] / page_height * img_height)
                ]
                
                # デバッグ: 座標変換の確認
                print(f"  [DEBUG JSON] Text: '{block['text'][:20]}...'")
                print(f"    PDF bbox: {bbox}")
                print(f"    Img bbox: {img_bbox}")
                print(f"    Img size: {img_width}x{img_height}, Page size: {page_width}x{page_height}")
                
                # 画像から実際のテキスト高さを測定
                if cv_img is not None:
                    refined_bbox = refine_bbox_vertical(cv_img, img_bbox, bg_color)
                    actual_text_height_px = refined_bbox[3] - refined_bbox[1]
                    print(f"    Refined bbox: {refined_bbox}, Text height: {actual_text_height_px}px")
                    
                    # フォントサイズは実際のテキスト高さから計算
                    final_font_size_pt = estimate_font_size_v5(actual_text_height_px, img_height, page_height)
                    print(f"    Calculated font size: {final_font_size_pt:.1f}pt")
                else:
                    final_font_size_pt = (bbox[3] - bbox[1]) * 0.7
                    print(f"    Fallback font size: {final_font_size_pt:.1f}pt")
                
                run.font.size = Pt(final_font_size_pt)
                
                # 色
                color_rgb = detect_text_color(cv_img, img_bbox, bg_color) if cv_img is not None else (0, 0, 0)
                run.font.color.rgb = RGBColor(color_rgb[0], color_rgb[1], color_rgb[2])
                print(f"    Detected color: {color_rgb}")
                
                # フォント・太字 (v5ロジック)
                if cv_img is not None:
                    ratio, v_var = detect_font_weight_v5(cv_img, img_bbox, bg_color)
                    selected_font, is_bold_detected = get_font_for_text_v5(ratio, v_var, color_rgb)
                    print(f"    Weight ratio: {ratio:.4f}, Stroke var: {v_var:.4f}")
                    print(f"    Selected font: {selected_font}, Bold: {is_bold_detected}")
                    
                    if "font_family" in block:
                        set_font_for_run(run, block["font_family"])
                    else:
                        run.font.name = selected_font
                    
                    if "is_bold" in block:
                        run.font.bold = block["is_bold"]
                    elif is_bold_detected:
                        run.font.bold = True
                else:
                    if "font_family" in block: set_font_for_run(run, block["font_family"])
                    if block.get("is_bold"): run.font.bold = True
    
    # 保存
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path

def main(pdf_path=None, output_path=None, json_path=None):
    """メイン処理 - 全ページ変換
    
    Args:
        pdf_path: PDFファイルパス（Noneの場合はデフォルト値を使用）
        output_path: 出力PPTXパス（Noneの場合はデフォルト値を使用）
        json_path: OCR補正JSONパス（Noneの場合はデフォルト値を使用）
    """
    global CURRENT_JSON_PATH
    
    # JSONパスを設定（CLIから渡された場合）
    if json_path:
        CURRENT_JSON_PATH = json_path
        print(f"Using JSON path: {CURRENT_JSON_PATH}")
    
    # デフォルト値
    if pdf_path is None:
        pdf_path = r"c:\Users\ishik\Documents\Antigravity\PDFパワポ化\webapp\temp_processing\129493ae-b15b-4089-9e12-07d76f8f3603\input.pdf"
    if output_path is None:
        final_output_dir = r"c:\Users\ishik\Documents\Antigravity\PDFパワポ化\output"
        os.makedirs(final_output_dir, exist_ok=True)
        output_path = os.path.join(final_output_dir, "webapp_test_output2.pptx")
    
    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    print(f"Total pages: {total_pages}")
    
    # 最初のページでサイズを取得
    first_page = doc[0]
    page_width = first_page.rect.width
    page_height = first_page.rect.height
    print(f"Page size: {page_width} x {page_height} points")
    
    # プレゼンテーションを作成（スライドサイズを設定）
    prs = Presentation()
    prs.slide_width = Emu(int(page_width * 12700))
    prs.slide_height = Emu(int(page_height * 12700))
    
    for page_num in range(total_pages):  # 全ページ処理
        global CURRENT_PAGE_NUM, OCR_ALL_WORDS
        CURRENT_PAGE_NUM = page_num + 1
        OCR_ALL_WORDS = []  # ページごとにリセット
        print(f"\n=== Processing page {page_num + 1}/{total_pages} ===")
        try:
            print(f"[{page_num+1}] Loading page...")
            page = doc[page_num]
            
            # ページ画像を生成
            print(f"[{page_num+1}] Generating page image...")
            page_image_path = os.path.join(OUTPUT_DIR, f"page_{page_num}.png")
            mat = fitz.Matrix(2.0, 2.0)  # Full resolution for quality
            pix = page.get_pixmap(matrix=mat)
            pix.save(page_image_path)
            print(f"[{page_num+1}] Page image saved.")
            
            # Create slide
            print(f"[{page_num+1}] Adding slide...")
            add_slide_from_page(prs, page, page_image_path, page_width, page_height)
            print(f"[{page_num+1}] Page completed successfully.")
            
            # LIGHTWEIGHT: メモリ解放とページ画像削除
            del pix
            import gc
            gc.collect()
            
            # ページ画像ファイルも削除してディスク容量節約
            try:
                if os.path.exists(page_image_path):
                    os.remove(page_image_path)
            except:
                pass
            
        except Exception as e:
            print(f"Error on page {page_num + 1}: {e}")
            import traceback
            traceback.print_exc()
            continue
    
    # 保存
    prs.save(output_path)
    print(f"\nFinal output: {output_path}")
    
    doc.close()
    print("Conversion complete!")
    return output_path

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="PDF to PPTX Converter - Precision Mode")
    parser.add_argument("--pdf", "-p", type=str, help="Input PDF file path")
    parser.add_argument("--output", "-o", type=str, help="Output PPTX file path")
    parser.add_argument("--json", "-j", type=str, help="OCR corrections JSON path")
    parser.add_argument("--log", "-l", type=str, help="Log file path")
    args = parser.parse_args()
    
    # ログ設定
    log_path = args.log if args.log else "debug_log_beans.txt"
    log_file = open(log_path, "w", encoding="utf-8", buffering=1)
    sys.stdout = log_file
    
    try:
        main(pdf_path=args.pdf, output_path=args.output, json_path=args.json)
    finally:
        log_file.close()

